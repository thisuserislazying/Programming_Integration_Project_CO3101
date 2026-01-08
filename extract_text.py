import os
import pymupdf 
from pptx import Presentation

# cau hinh thu muc va file dau ra
RAW_FOLDER = "raw_files"
BACKUP_FILE = "Backup/KTCT.txt"

def get_all_files(source_input):
    # lay danh sach tat ca cac file trong thu muc
    file_list = []
    if isinstance(source_input, list):
        for item in source_input:
            file_list.extend(get_all_files(item))
    elif os.path.isdir(source_input):
        for root, _, files in os.walk(source_input):
            for file in files:
                if not file.startswith("~$"): 
                    file_list.append(os.path.join(root, file))
    elif os.path.isfile(source_input):
        file_list.append(source_input)
    return file_list

def process_pdf_content(file_path):
    # xu ly noi dung file pdf
    try:
        doc = pymupdf.open(file_path)
        text_out = ""
        for page in doc:
            blocks = page.get_text("blocks")
            for b in blocks:
                block_text = b[4].strip()
                if block_text.isdigit() and len(block_text) < 4:
                    continue
                text_out += block_text + "\n"
        return text_out
    except Exception:
        return ""

def process_pptx_content(file_path):
    # xu ly noi dung file pptx
    try:
        prs = Presentation(file_path)
        text_out = []
        for slide in prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for paragraph in shape.text_frame.paragraphs:
                        p_text = ''.join(run.text for run in paragraph.runs)
                        if p_text.strip():
                            slide_text.append(p_text.strip())
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = [cell.text_frame.text.strip() for cell in row.cells if cell.text_frame.text.strip()]
                        if row_text:
                            slide_text.append(" | ".join(row_text))
            
            if slide_text:
                text_out.append("\n".join(slide_text))
        return "\n\n".join(text_out)
    except Exception:
        return ""

def smart_extract(raw_source, backup_file):
    # tao thu muc neu chua co
    os.makedirs(os.path.dirname(backup_file), exist_ok=True)
    all_files = get_all_files(raw_source)

    # mo file ghi du lieu
    with open(backup_file, "w", encoding="utf-8") as f:
        for file_path in all_files:
            content = ""
            ext = os.path.splitext(file_path)[1].lower()
            
            # phan loai file de xu ly
            if ext == ".pdf":
                content = process_pdf_content(file_path)
            elif ext in [".pptx", ".ppt"]:
                content = process_pptx_content(file_path)
            
            # ghi noi dung vao file tong
            if content:
                f.write(f"\n{'='*20}\nSOURCE_START: {os.path.basename(file_path)}\n{'='*20}\n")
                f.write(content)
                f.write(f"\nSOURCE_END: {os.path.basename(file_path)}\n")
    
    # thong bao hoan tat
    print(f"[*] da xu ly xong extract text: {backup_file}")

if __name__ == "__main__":
    if os.path.exists(RAW_FOLDER):
        smart_extract(RAW_FOLDER, BACKUP_FILE)